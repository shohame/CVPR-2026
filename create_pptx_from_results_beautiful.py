"""
Create a polished PowerPoint presentation from an HTML results table.
Each row becomes one slide:
1. Title and instrument name at the top
2. Instrument image below the header
3. Spectrogram graphs
4. Audio playback under each graph
"""

import mimetypes
import subprocess
from pathlib import Path

from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# Paths
results_dir = Path(r"d:\work\python\CVPR-2026\hearing_the_shape_of_the_drum-main\html_assets\results")
html_file = results_dir / "results.html"
output_pptx = results_dir / "Recovered_Audio_Results_beautiful.pptx"


# Slide constants
SLIDE_W = 13.0
SLIDE_H = 7.5

# Theme
COL_BG = RGBColor(246, 248, 252)
COL_NAVY = RGBColor(18, 31, 55)
COL_NAVY_2 = RGBColor(37, 53, 84)
COL_TEXT = RGBColor(18, 24, 38)
COL_MUTED = RGBColor(94, 105, 124)
COL_WHITE = RGBColor(255, 255, 255)
COL_ACCENT = RGBColor(32, 176, 158)
COL_CARD = RGBColor(255, 255, 255)
COL_LINE = RGBColor(218, 224, 235)
COL_OURS = RGBColor(10, 130, 110)


def set_slide_background(slide, color=COL_BG):
    """Set slide background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False,
                color=COL_TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    """Add formatted text."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = valign

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_shape(slide, shape_type, left, top, width, height, fill_color,
              line_color=None, line_width=0):
    """Add a simple filled shape."""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None or line_width == 0:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def svg_to_png(svg_path):
    """Convert an SVG file to PNG using ImageMagick or Inkscape."""
    svg_path = Path(svg_path)
    png_path = svg_path.with_suffix(".png")
    if png_path.exists():
        return png_path

    commands = [
        ["magick", str(svg_path), str(png_path)],
        ["convert", str(svg_path), str(png_path)],
        ["inkscape", str(svg_path), "--export-type=png", f"--export-filename={png_path}"],
    ]

    for cmd in commands:
        try:
            subprocess.run(cmd, check=True, capture_output=True, timeout=20)
            if png_path.exists():
                return png_path
        except Exception:
            pass

    print(f"  Warning: Could not convert SVG: {svg_path}")
    return None


def prepare_image_path(image_path):
    """Convert SVG to PNG when needed and return a usable image path."""
    image_path = Path(image_path)
    if image_path.suffix.lower() == ".svg":
        converted = svg_to_png(image_path)
        return converted if converted else None
    return image_path


def add_picture_contain(slide, image_path, left, top, width, height):
    """Add an image without distortion, fitted inside a bounding box."""
    image_path = prepare_image_path(image_path)
    if not image_path or not Path(image_path).exists():
        return None

    try:
        with Image.open(image_path) as im:
            img_w, img_h = im.size
    except Exception:
        return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))

    box_ratio = width / height
    img_ratio = img_w / img_h

    if img_ratio > box_ratio:
        final_w = width
        final_h = width / img_ratio
    else:
        final_h = height
        final_w = height * img_ratio

    final_left = left + (width - final_w) / 2
    final_top = top + (height - final_h) / 2

    return slide.shapes.add_picture(
        str(image_path),
        Inches(final_left),
        Inches(final_top),
        width=Inches(final_w),
        height=Inches(final_h),
    )


def add_top_header(slide, setup_name):
    """Add a top header with deck title and instrument name."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.82, COL_NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0.82, SLIDE_W, 0.05, COL_ACCENT)

    add_textbox(
        slide,
        "Recovered Audio Results",
        0.45,
        0.16,
        5.6,
        0.42,
        font_size=23,
        bold=True,
        color=COL_WHITE,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        setup_name.upper(),
        6.1,
        0.15,
        6.45,
        0.45,
        font_size=21,
        bold=True,
        color=COL_WHITE,
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_title_slide(prs):
    """Create the opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COL_NAVY)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 6.95, SLIDE_W, 0.08, COL_ACCENT)
    add_textbox(
        slide,
        "Recovered Audio Results",
        0.8,
        2.25,
        11.4,
        0.8,
        font_size=40,
        bold=True,
        color=COL_WHITE,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "Comparison of audio recovery methods from vibrating surfaces",
        1.25,
        3.15,
        10.5,
        0.45,
        font_size=19,
        bold=False,
        color=RGBColor(220, 228, 241),
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "Figure 5 results",
        4.85,
        4.05,
        3.3,
        0.42,
        font_size=16,
        bold=True,
        color=COL_ACCENT,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_audio_control(slide, audio_path, left, top, width, height, label="PLAY SOUND"):
    """Embed audio when possible. Fallback creates a linked button."""
    audio_path = Path(audio_path)
    if not audio_path.exists():
        add_textbox(
            slide,
            "Audio missing",
            left,
            top,
            width,
            height,
            font_size=8,
            bold=True,
            color=COL_MUTED,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        return

    mime_type, _ = mimetypes.guess_type(str(audio_path))
    mime_type = mime_type or "audio/wav"

    try:
        slide.shapes.add_movie(
            str(audio_path),
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            mime_type=mime_type,
        )
        return
    except Exception as e:
        print(f"   Warning: Could not embed audio ({e}). Creating linked button instead.")

    button = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        COL_ACCENT,
        line_color=COL_ACCENT,
        line_width=0.5,
    )
    button.click_action.hyperlink.address = str(audio_path)
    tf = button.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COL_WHITE


def add_method_panel(slide, method, idx, total):
    """Add spectrogram, method label, and audio for one method."""
    margin_x = 0.34
    gap = 0.13
    panel_top = 3.18
    panel_h = 3.9
    panel_w = (SLIDE_W - 2 * margin_x - gap * (total - 1)) / total
    left = margin_x + idx * (panel_w + gap)

    fill = COL_CARD
    line = COL_ACCENT if method["name"] == "Ours" else COL_LINE
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, panel_top, panel_w, panel_h, fill, line, 1.0)

    label_color = COL_OURS if method["name"] == "Ours" else COL_TEXT
    add_textbox(
        slide,
        method["name"],
        left + 0.08,
        panel_top + 0.11,
        panel_w - 0.16,
        0.28,
        font_size=12,
        bold=True,
        color=label_color,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    spec_top = panel_top + 0.48
    spec_h = 2.55
    spec_w = panel_w - 0.18
    try:
        add_picture_contain(slide, method["spec_path"], left + 0.09, spec_top, spec_w, spec_h)
        print(f"   Added {method['name']} spectrogram")
    except Exception as e:
        print(f"   Warning: Could not add {method['name']} spectrogram - {e}")

    add_audio_control(
        slide,
        method["audio_path"],
        left + 0.28,
        panel_top + 3.28,
        panel_w - 0.56,
        0.38,
    )


def extract_rows(html_path):
    """Read the results HTML and extract table rows."""
    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    table = soup.find("table", class_="grid")
    if table is None:
        raise RuntimeError("Could not find table with class 'grid' in results.html")

    rows = table.find_all("tr")
    return rows[1:]


def build_presentation():
    data_rows = extract_rows(html_file)
    print(f"Found {len(data_rows)} data rows")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs)

    methods = ["Single Point", "Average", "Delay & Sum", "Ours", "Input"]

    for row_idx, row in enumerate(data_rows):
        print(f"\nProcessing row {row_idx + 1}...")

        cells = row.find_all("td")
        if len(cells) < 6:
            print(f"  Skipping row {row_idx + 1} - not enough cells")
            continue

        setup_cell = cells[0]
        setup_img_tag = setup_cell.find("img")
        setup_caption_tag = setup_cell.find("figcaption")

        if not setup_img_tag or not setup_caption_tag:
            print(f"  Skipping row {row_idx + 1} - missing setup info")
            continue

        setup_img_path = results_dir / setup_img_tag.get("src")
        setup_name = setup_caption_tag.get_text(strip=True)
        print(f"  Setup: {setup_name}")

        method_data = []
        for method_idx in range(1, 6):
            cell = cells[method_idx]
            spec_img_tag = cell.find("img", class_="spec")
            audio_tag = cell.find("audio")

            if spec_img_tag and audio_tag:
                method_data.append({
                    "name": methods[method_idx - 1],
                    "spec_path": results_dir / spec_img_tag.get("src"),
                    "audio_path": results_dir / audio_tag.get("src"),
                })

        if len(method_data) < 5:
            print("  Skipping row - missing method data")
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        add_top_header(slide, setup_name)

        # Instrument image section
        add_textbox(
            slide,
            "Instrument / vibrating object",
            0.6,
            0.98,
            3.0,
            0.3,
            font_size=11,
            bold=True,
            color=COL_MUTED,
            align=PP_ALIGN.LEFT,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.65, 1.0, 3.7, 1.65, COL_WHITE, COL_LINE, 1)
        try:
            add_picture_contain(slide, setup_img_path, 4.85, 1.12, 3.3, 1.38)
            print("   Added setup image")
        except Exception as e:
            print(f"   Warning: Could not add setup image - {e}")

        add_textbox(
            slide,
            "Spectrogram comparison and recovered sound",
            0.55,
            2.72,
            12.0,
            0.32,
            font_size=14,
            bold=True,
            color=COL_TEXT,
            align=PP_ALIGN.LEFT,
            valign=MSO_ANCHOR.MIDDLE,
        )

        for method_idx, method in enumerate(method_data):
            add_method_panel(slide, method, method_idx, len(method_data))

    prs.save(str(output_pptx))
    print(f"\nSuccessfully created: {output_pptx}")


if __name__ == "__main__":
    build_presentation()
