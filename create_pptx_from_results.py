"""
Create PowerPoint presentation from HTML results table
Each row becomes a slide with setup image, spectrograms, and embedded audio files
"""

import os
import re
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import subprocess

# Paths
results_dir = Path(r"d:\work\python\CVPR-2026\hearing_the_shape_of_the_drum-main\html_assets\results")
html_file = results_dir / "results.html"
output_pptx = results_dir / "Recovered_Audio_Results_centered.pptx"

# Parse HTML
with open(html_file, 'r', encoding='utf-8') as f:
    soup = BeautifulSoup(f.read(), 'html.parser')

# Extract table rows
table = soup.find('table', class_='grid')
rows = table.find_all('tr')

# Skip header row
data_rows = rows[1:]

print(f"Found {len(data_rows)} data rows")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13)
prs.slide_height = Inches(7.5)

# Add title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Recovered Audio Results (Figure 5)"
subtitle.text = "Comparison of audio recovery methods from vibrating surfaces"

# Helper function to convert SVG to PNG
def svg_to_png(svg_path, png_path):
    """Convert SVG to PNG using Inkscape or ImageMagick"""
    if Path(png_path).exists():
        return png_path
    
    try:
        # Try using ImageMagick (convert)
        subprocess.run(['convert', str(svg_path), str(png_path)], 
                      check=True, capture_output=True, timeout=10)
        return png_path
    except:
        try:
            # Try using Inkscape
            subprocess.run(['inkscape', '-l', str(png_path), str(svg_path)],
                          check=True, capture_output=True, timeout=10)
            return png_path
        except:
            print(f"  Warning: Could not convert {svg_path}")
            return None

# Process each row
for row_idx, row in enumerate(data_rows):
    print(f"\nProcessing row {row_idx + 1}...")
    
    cells = row.find_all('td')
    if len(cells) < 6:
        print(f"  Skipping row {row_idx + 1} - not enough cells")
        continue
    
    # Extract setup info
    setup_cell = cells[0]
    setup_img_tag = setup_cell.find('img')
    setup_caption_tag = setup_cell.find('figcaption')
    
    if not setup_img_tag or not setup_caption_tag:
        print(f"  Skipping row {row_idx + 1} - missing setup info")
        continue
    
    setup_img_path = results_dir / setup_img_tag.get('src')
    setup_name = setup_caption_tag.get_text(strip=True)
    
    print(f"  Setup: {setup_name}")
    
    # Extract method data (5 columns: Single, Average, Delay&Sum, Ours, Input)
    methods = ['Single Point', 'Average', 'Delay & Sum', 'Ours', 'Input']
    method_data = []
    
    for method_idx in range(1, 6):
        cell = cells[method_idx]
        spec_img_tag = cell.find('img', class_='spec')
        audio_tag = cell.find('audio')
        
        if spec_img_tag and audio_tag:
            spec_path = results_dir / spec_img_tag.get('src')
            audio_path = results_dir / audio_tag.get('src')
            method_data.append({
                'name': methods[method_idx - 1],
                'spec_path': spec_path,
                'audio_path': audio_path
            })
    
    if len(method_data) < 5:
        print(f"  Skipping row - missing method data")
        continue
    
    # Create slide for this setup
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Main title in black
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.2), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Recovered Audio Results"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # Instrument image centered under the title
    try:
        slide.shapes.add_picture(
            str(setup_img_path),
            Inches(5.2),
            Inches(0.9),
            width=Inches(2.6),
            height=Inches(1.5),
        )
        print("   Added setup image")
    except Exception as e:
        print(f"   Warning: Could not add setup image - {e}")

    # Large instrument name centered below the image
    instrument_box = slide.shapes.add_textbox(Inches(0.4), Inches(2.45), Inches(12.2), Inches(0.5))
    instrument_frame = instrument_box.text_frame
    instrument_frame.word_wrap = True
    p = instrument_frame.paragraphs[0]
    p.text = setup_name.upper()
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # Large spectrogram grid that covers most of the slide width
    spec_width = Inches(2.4)
    spec_height = Inches(2.7)
    col_gap = Inches(0.08)
    start_left = Inches(0.34)
    start_top = Inches(3.0)
    
    for method_idx, method in enumerate(method_data):
        left = start_left + (method_idx * (spec_width + col_gap))
        top = start_top
        
        # Add spectrogram
        spec_path = method['spec_path']
        if spec_path.suffix.lower() == '.svg':
            # Convert SVG to PNG
            png_path = spec_path.with_suffix('.png')
            converted_path = svg_to_png(str(spec_path), str(png_path))
            if converted_path:
                spec_path = Path(converted_path)
            else:
                # Skip if conversion failed
                continue
        
        try:
            slide.shapes.add_picture(str(spec_path), left, top, width=spec_width)
            print(f"   Added {method['name']} spectrogram")
        except Exception as e:
            print(f"   Warning: Could not add {method['name']} spectrogram - {e}")
        
        # Add method label
        label_box = slide.shapes.add_textbox(left, top + spec_height + Inches(0.08),
                            spec_width, Inches(0.35))
        label_frame = label_box.text_frame
        label_frame.word_wrap = True
        p = label_frame.paragraphs[0]
        p.text = method['name']
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        
        # Add embedded audio player
        audio_path = method['audio_path']
        if Path(audio_path).exists():
            try:
                # Add embedded audio/movie object
                # This creates a clickable media player on the slide
                audio_left = left + Inches(0.35)
                audio_top = top + spec_height + Inches(0.45)
                audio_width = Inches(1.7)
                audio_height = Inches(0.42)
                
                # Add movie/audio shape with embedded file
                movie = slide.shapes.add_movie(
                    str(audio_path),
                    audio_left,
                    audio_top,
                    audio_width,
                    audio_height
                )
                
                # Style the movie player
                # The movie object represents the embedded media player
                print(f"   Embedded {method['name']} audio file for playback")
                
            except Exception as e:
                print(f"   Warning: Could not embed audio - {e}, creating button instead")
                # Fallback: create a clickable button
                try:
                    audio_button = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        left + Inches(0.35),
                        top + spec_height + Inches(0.45),
                        Inches(1.7),
                        Inches(0.42)
                    )
                    audio_button.fill.solid()
                    audio_button.fill.fore_color.rgb = RGBColor(32, 238, 193)
                    audio_button.line.color.rgb = RGBColor(255, 255, 255)
                    audio_button.line.width = Pt(1)
                    
                    text_frame = audio_button.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = False
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = text_frame.paragraphs[0]
                    p.text = "🔊 Play"
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(18, 18, 18)
                    p.alignment = PP_ALIGN.CENTER
                except:
                    pass
        else:
            print(f"   Warning: Audio file not found at {audio_path}")

# Save presentation
try:
    prs.save(str(output_pptx))
    print(f"\n✅ Successfully created: {output_pptx}")
except Exception as e:
    print(f"\n❌ Error saving presentation: {e}")
